# -*- coding: utf-8 -*-
"""
Сборка презентации (python-pptx, т.к. Node.js/pptxgenjs на этой машине нет).
Тёмная тема в цвет уже готовых графиков анализа.

Запуск: python build_presentation.py
"""

from itertools import count

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Палитра и константы
# ---------------------------------------------------------------------------

BG = RGBColor(0x10, 0x20, 0x3C)
CARD = RGBColor(0x1C, 0x2E, 0x4A)
TERRACOTTA = RGBColor(0xC1, 0x44, 0x0E)
CREAM = RGBColor(0xE8, 0xE6, 0xDF)
MUTED = RGBColor(0x9B, 0x97, 0x8D)
GOLD = RGBColor(0xE0, 0xA4, 0x58)
SAGE = RGBColor(0x5B, 0x8A, 0x72)
SALMON = RGBColor(0xE0, 0x7A, 0x5F)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)


def img_size(path):
    with Image.open(path) as im:
        return im.size  # (w, h) px


def fit_box(path, max_w, max_h):
    """Возвращает (w, h) в EMU, вписывая изображение в max_w x max_h с сохранением пропорций."""
    pw, ph = img_size(path)
    aspect = pw / ph
    w, h = max_w, max_w / aspect
    if h > max_h:
        h = max_h
        w = max_h * aspect
    return w, h


def new_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def set_no_autosize(tf):
    tf.word_wrap = True
    try:
        el = tf._txBody
        bodyPr = el.find(qn('a:bodyPr'))
        for tag in ('a:normAutofit', 'a:spAutoFit'):
            e = bodyPr.find(qn(tag))
            if e is not None:
                bodyPr.remove(e)
    except Exception:
        pass


def add_text(slide, x, y, w, h, text, size=14, color=CREAM, bold=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.0, space_after=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    set_no_autosize(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return box


def add_rich(slide, x, y, w, h, segments, size=14, font=BODY_FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """segments: список параграфов, каждый параграф -- список (text, color, bold, italic, size?)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    set_no_autosize(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for seg in para:
            t, c = seg[0], seg[1]
            b = seg[2] if len(seg) > 2 else False
            it = seg[3] if len(seg) > 3 else False
            sz = seg[4] if len(seg) > 4 else size
            run = p.add_run()
            run.text = t
            run.font.size = Pt(sz)
            run.font.color.rgb = c
            run.font.bold = b
            run.font.italic = it
            run.font.name = font
    return box


def add_title(slide, text, subtitle=None, color=CREAM):
    add_text(slide, MARGIN, Inches(0.42), SLIDE_W - 2 * MARGIN, Inches(0.7),
              text, size=28, color=color, bold=True, font=HEAD_FONT)
    if subtitle:
        add_text(slide, MARGIN, Inches(1.02), SLIDE_W - 2 * MARGIN, Inches(0.35),
                  subtitle, size=13, color=MUTED, font=BODY_FONT, italic=True)


def add_circle_num(slide, x, y, d, number, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BG
    run.font.name = HEAD_FONT
    return shp


def add_image_fit(slide, path, x, y, max_w, max_h, center_h=True, center_v=False):
    w, h = fit_box(path, max_w, max_h)
    ox = x + (max_w - w) / 2 if center_h else x
    oy = y + (max_h - h) / 2 if center_v else y
    slide.shapes.add_picture(path, ox, oy, width=Emu(int(w)), height=Emu(int(h)))
    return ox, oy, w, h


def add_card(slide, x, y, w, h, color=CARD):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.04
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def footer_page(slide, n):
    add_text(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.4), Inches(0.5), Inches(0.3),
              str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Презентация
# ---------------------------------------------------------------------------

prs = new_pres()
page_num = count(2)

# ---- Слайд 1: титул ------------------------------------------------------
s = add_slide(prs)
img_w = Inches(5.7)
ox, oy, w, h = add_image_fit(s, "S2_RGB_Barringer__Meteor_Crater______.png",
                              SLIDE_W - Inches(0.5) - img_w, Inches(0.5), img_w, SLIDE_H - Inches(1.0),
                              center_h=True, center_v=True)
add_text(s, Inches(0.7), Inches(1.7), Inches(6.2), Inches(2.6),
          "Дистанционное зондирование\nземных кратеров как инструмент\nраспознавания импактных структур",
          size=26, color=CREAM, bold=True, font=HEAD_FONT, line_spacing=1.12)
add_text(s, Inches(0.7), Inches(4.15), Inches(6.2), Inches(0.7),
          "Морфометрия по DEM, Sentinel-1 и Sentinel-2:\nBarringer, Lonar, Wolfe Creek",
          size=14.5, color=GOLD, italic=True, line_spacing=1.15)
add_text(s, Inches(0.7), Inches(5.6), Inches(6.2), Inches(0.4),
          "Баканас Е.С.¹˒², Бадак Л.А.²", size=15, color=CREAM, bold=True)
add_text(s, Inches(0.7), Inches(6.0), Inches(6.2), Inches(1.25),
          "¹ Научный центр оперативного мониторинга Земли АО «Российские космические системы», Москва\n"
          "² Институт астрономии Российской академии наук, Москва",
          size=10.5, color=MUTED, line_spacing=1.25)

# ---- Слайд 2: актуальность -----------------------------------------------
s = add_slide(prs)
add_title(s, "Зачем изучать ударные кратеры")
items2 = [
    ("Реконструкция геологической истории Земли", TERRACOTTA),
    ("Оценка роли импактных событий в изменении природной среды", GOLD),
    ("Изучение свойств астероидов и других малых тел", SAGE),
    ("Совершенствование методов поиска потенциально опасных объектов", SALMON),
]
y0 = Inches(1.75)
row_h = Inches(1.02)
for i, (txt, color) in enumerate(items2):
    y = y0 + i * row_h
    add_circle_num(s, MARGIN, y, Inches(0.55), i + 1, color)
    add_text(s, MARGIN + Inches(0.85), y - Inches(0.02), Inches(10.8), Inches(0.65),
              txt, size=17, color=CREAM, anchor=MSO_ANCHOR.MIDDLE)

add_card(s, MARGIN, Inches(5.95), SLIDE_W - 2 * MARGIN, Inches(1.15), color=CARD)
add_text(s, MARGIN + Inches(0.35), Inches(6.12), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(0.85),
          "Земной реестр импактных структур неполон: многие кратеры разрушены эрозией и тектоникой, "
          "перекрыты осадками, растительностью или водой — поэтому выявление и картографирование "
          "сохранившихся структур по данным ДЗЗ имеет существенное научное и прикладное значение.",
          size=13, color=CREAM, italic=True, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
footer_page(s, next(page_num))

# ---- Слайд 3: как распознают -----------------------------------------------
s = add_slide(prs)
add_title(s, "Двухэтапное распознавание")
col_w = Inches(5.75)
gap = Inches(0.35)
x1 = MARGIN
x2 = MARGIN + col_w + gap
y_top = Inches(1.75)
col_h = Inches(3.9)

add_card(s, x1, y_top, col_w, col_h, color=CARD)
add_circle_num(s, x1 + Inches(0.35), y_top + Inches(0.35), Inches(0.55), 1, TERRACOTTA)
add_text(s, x1 + Inches(1.1), y_top + Inches(0.38), col_w - Inches(1.4), Inches(0.5),
          "Морфология (дистанционно)", size=18, color=TERRACOTTA, bold=True, font=HEAD_FONT)
morf_items = ["округлая или близкая к округлой форма", "вал и внутренний склон", "террасы",
              "центральное поднятие или понижение", "зона выбросов"]
add_rich(s, x1 + Inches(0.5), y_top + Inches(1.25), col_w - Inches(1.0), Inches(2.4),
          [[("•  ", TERRACOTTA, True), (t, CREAM)] for t in morf_items],
          size=14.5, line_spacing=1.55)

add_card(s, x2, y_top, col_w, col_h, color=CARD)
add_circle_num(s, x2 + Inches(0.35), y_top + Inches(0.35), Inches(0.55), 2, GOLD)
add_text(s, x2 + Inches(1.1), y_top + Inches(0.38), col_w - Inches(1.4), Inches(0.5),
          "Геология (подтверждение)", size=18, color=GOLD, bold=True, font=HEAD_FONT)
geo_items = ["брекчии, конусы разрушения", "деформации минералов",
             "высокобарные полиморфные модификации", "состав и возраст пород",
             "структура разреза, распределение выбросов"]
add_rich(s, x2 + Inches(0.5), y_top + Inches(1.25), col_w - Inches(1.0), Inches(2.4),
          [[("•  ", GOLD, True), (t, CREAM)] for t in geo_items],
          size=14.5, line_spacing=1.55)

add_text(s, MARGIN, Inches(5.9), SLIDE_W - 2 * MARGIN, Inches(1.0),
          "Сходные формы бывают вулканического, тектонического, карстового или эрозионного "
          "происхождения — морфология формирует рабочую гипотезу, а окончательная идентификация "
          "требует сопоставления с независимыми геологическими данными.",
          size=13.5, color=MUTED, italic=True, line_spacing=1.25)
footer_page(s, next(page_num))

# ---- Слайд 3b: почему одной морфологии недостаточно -----------------------------------------------
s = add_slide(prs)
add_title(s, "Почему одной морфологии недостаточно",
           "И зачем тогда вообще нужен дистанционный этап")

tiers = [
    ("1", "Морфология и ДЗЗ (наш этап)", "весь земной шар, часы, бесплатно",
     "не доказательство: похожую форму даёт вулканизм и эрозия", TERRACOTTA, Inches(11.8)),
    ("2", "Геофизика", "сужает круг кандидатов",
     "тоже не доказательство: аномалию дают и другие тела", GOLD, Inches(9.3)),
    ("3", "Лаборатория (петрография + геохимия)", "шлифы под микроскопом, месяцы работы",
     "PDF, шаттер-конусы, коэсит — доказательство (Koeberl, 2004)", SAGE, Inches(6.8)),
]
y_t = Inches(1.65)
bar_h = Inches(1.2)
gap_t = Inches(0.15)
for num, head, cheap, strict, color, bar_w in tiers:
    x_t = MARGIN + (SLIDE_W - 2 * MARGIN - bar_w) / 2
    add_card(s, x_t, y_t, bar_w, bar_h, color=CARD)
    add_circle_num(s, x_t + Inches(0.25), y_t + Inches(0.3), Inches(0.55), num, color)
    add_text(s, x_t + Inches(1.0), y_t + Inches(0.14), bar_w - Inches(1.3), Inches(0.35),
              head, size=15, color=color, bold=True, font=HEAD_FONT)
    add_rich(s, x_t + Inches(1.0), y_t + Inches(0.56), bar_w - Inches(1.3), Inches(0.58),
              [[("+  ", SAGE, True), (cheap, CREAM, False, False, 11)],
               [("−  ", SALMON, True), (strict, MUTED, False, True, 10.5)]],
              line_spacing=1.15)
    y_t += bar_h + gap_t

add_card(s, MARGIN, y_t + Inches(0.05), SLIDE_W - 2 * MARGIN, Inches(1.25), color=CARD)
add_rich(s, MARGIN + Inches(0.35), y_t + Inches(0.24), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(1.0),
          [[("Стоит ли тогда заниматься дистанционным этапом? ", GOLD, True, False, 13.5),
            ("Да — лабораторная проверка дорога и медленна и не может применяться подряд ко всем "
             "потенциальным структурам. Наш метод — быстрый воспроизводимый фильтр, который сужает "
             "круг кандидатов для дорогой проверки, а не заменяет её.", CREAM, False, True, 13.5)]],
          line_spacing=1.28)
footer_page(s, next(page_num))

# ---- Слайд 4: данные ДЗЗ -----------------------------------------------
s = add_slide(prs)
add_title(s, "Три источника данных ДЗЗ, три роли")
cols4 = [
    ("Оптика / ИК", "Sentinel-2", TERRACOTTA,
     "состав и степень выветрелости пород, зоны выбросов, нарушенный грунт"),
    ("Цифровая модель рельефа", "SRTM / GLO-30 / 3DEP", GOLD,
     "форма вала, глубина чаши, крутизна склонов, террасы, степень деградации"),
    ("Радар", "Sentinel-1", SAGE,
     "шероховатость и текстура поверхности; работает при облачности и без солнечного освещения"),
]
col_w4 = Inches(3.85)
gap4 = Inches(0.29)
y4 = Inches(1.8)
h4 = Inches(3.55)
for i, (head, sub, color, desc) in enumerate(cols4):
    x = MARGIN + i * (col_w4 + gap4)
    add_card(s, x, y4, col_w4, h4, color=CARD)
    add_circle_num(s, x + Inches(0.3), y4 + Inches(0.3), Inches(0.5), i + 1, color)
    add_text(s, x + Inches(0.3), y4 + Inches(1.0), col_w4 - Inches(0.6), Inches(0.5),
              head, size=16.5, color=color, bold=True, font=HEAD_FONT)
    add_text(s, x + Inches(0.3), y4 + Inches(1.48), col_w4 - Inches(0.6), Inches(0.35),
              sub, size=11.5, color=MUTED, italic=True)
    add_text(s, x + Inches(0.3), y4 + Inches(2.0), col_w4 - Inches(0.6), Inches(1.4),
              desc, size=13, color=CREAM, line_spacing=1.3)

add_text(s, MARGIN, Inches(5.65), SLIDE_W - 2 * MARGIN, Inches(1.15),
          "Развитие глобальных радарных моделей рельефа (включая TanDEM-X) создало возможности "
          "для систематического морфометрического анализа известных структур и поиска новых "
          "кандидатов (Gottwald et al., 2017).",
          size=12.5, color=MUTED, italic=True, line_spacing=1.25)
footer_page(s, next(page_num))

# ---- Слайд 5: метод/конвейер -----------------------------------------------
s = add_slide(prs)
add_title(s, "Наш воспроизводимый конвейер",
           "Один и тот же скрипт, применённый без изменений к трём разным кратерам")
steps = [
    ("Координаты", "справочник кратеров или клик по центру на снимке", TERRACOTTA),
    ("N азимутов", "радиальные лучи вокруг центра кратера", SALMON),
    ("Выборка вдоль лучей", "DEM + VV/VH (Sentinel-1) + slope + hillshade\n(Google Earth Engine)", GOLD),
    ("CSV → Python", "дно, гребень вала, диаметр, глубина,\nнеоднородность радиуса вала (CV)", SAGE),
    ("Сопоставление", "с публикациями и эталонной кривой\nглубина/диаметр", RGBColor(0x8a, 0x6b, 0xb0)),
]
n = len(steps)
x0 = MARGIN
avail_in = SLIDE_W.inches - 2 * MARGIN.inches
step_w = Inches((avail_in - 0.35 * (n - 1)) / n)
gap5 = Inches(0.35)
y5 = Inches(2.2)
d5 = Inches(0.75)
for i, (head, desc, color) in enumerate(steps):
    x = x0 + i * (step_w + gap5)
    add_circle_num(s, x + (step_w - d5) / 2, y5, d5, i + 1, color)
    if i < n - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_w - Inches(0.05), y5 + Inches(0.27),
                                    gap5 + Inches(0.1), Inches(0.22))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()
        arrow.shadow.inherit = False
    add_text(s, x, y5 + Inches(1.05), step_w, Inches(0.55),
              head, size=14.5, color=color, bold=True, font=HEAD_FONT,
              align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_text(s, x, y5 + Inches(1.65), step_w, Inches(1.6),
              desc, size=11.5, color=CREAM, align=PP_ALIGN.CENTER, line_spacing=1.2)

add_card(s, MARGIN, Inches(5.9), SLIDE_W - 2 * MARGIN, Inches(1.15), color=CARD)
add_text(s, MARGIN + Inches(0.35), Inches(6.08), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(0.85),
          "Источник рельефа выбирается автоматически по приоритету: лидар USGS 3DEP (1 м, США) → "
          "Copernicus GLO-30 (~30 м, глобально) → SRTM (запасной вариант).",
          size=12.5, color=CREAM, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
footer_page(s, next(page_num))


# ---- Слайды 6-8: кейсы -----------------------------------------------
def case_slide(title_txt, chart_path, s2_path, color, stat_rows, note):
    s = add_slide(prs)
    add_title(s, title_txt)

    left_x = MARGIN
    left_w = Inches(7.85)
    top_y = Inches(1.55)
    add_image_fit(s, chart_path, left_x, top_y, left_w, Inches(5.55),
                  center_h=False, center_v=False)

    right_x = MARGIN + left_w + Inches(0.3)
    right_w = SLIDE_W - MARGIN - right_x

    thumb_h = Inches(2.35)
    add_image_fit(s, s2_path, right_x, top_y, right_w, thumb_h, center_h=True, center_v=False)

    ty = top_y + thumb_h + Inches(0.15)
    card_h = Inches(1.55) + Inches(0.42) * len(stat_rows)
    add_card(s, right_x, ty, right_w, card_h, color=CARD)

    yy = ty + Inches(0.25)
    pad = Inches(0.3)
    for label, our, pub in stat_rows:
        add_rich(s, right_x + pad, yy, right_w - 2 * pad, Inches(0.36),
                  [[(label + "  ", MUTED, False, False, 12), (our, color, True, False, 16),
                    ("   публ. " + pub, MUTED, False, True, 11)]])
        yy += Inches(0.42)
    yy += Inches(0.1)
    add_text(s, right_x + pad, yy, right_w - 2 * pad, card_h - (yy - ty) - Inches(0.2),
              note, size=11, color=CREAM, italic=True, line_spacing=1.22)
    footer_page(s, next(page_num))
    return s


case_slide("Кейс 1 — Барринджер (Meteor Crater)",
           "case_Barringer.png", "S2_RGB_Barringer__Meteor_Crater______.png", TERRACOTTA,
           [("Диаметр:", "1200 м", "1200 м"),
            ("Глубина:", "162 м", "170 м"),
            ("d/D:", "0.135", "0.142")],
           "Осадочные породы, ~50 тыс. лет. Наилучшее совпадение с публикацией — "
           "валидация метода на эталонном объекте.")

case_slide("Кейс 2 — Lonar",
           "case_Lonar.png", "S2_RGB_Lonar_______.png", SAGE,
           [("Диаметр:", "1890 м", "1830 м"),
            ("Глубина:", "111 м", "~137 м"),
            ("d/D:", "0.059", "0.075")],
           "Базальт (Деканские траппы). Два конфликтующих датирования: ~570 или ~37.5 тыс. лет. "
           "Глубина занижена — вероятно, DEM видит поверхность озера в кратере, а не истинное дно.")

case_slide("Кейс 3 — Wolfe Creek",
           "case_WolfeCreek.png", "S2_RGB_Wolfe_Creek___________.png", GOLD,
           [("Диаметр:", "930 м", "892 м"),
            ("Глубина:", "46 м", "178 м*")],
           "Песчаник, ~120 тыс. лет. *Исходная глубина по публикации — кратер занесён ~120 м песка; "
           "наш профиль отражает современную, а не исходную форму.")

# ---- Слайд 9: радар -----------------------------------------------
s = add_slide(prs)
add_title(s, "Радар видит вал независимо от освещения и облачности",
           "Sentinel-1, среднее VV на валу против дна кратера")

s1_imgs = [
    ("S1_VVVH_Barringer__Meteor_Crater______.png", "Барринджер", "+6.9 дБ", TERRACOTTA),
    ("S1_VVVH_Lonar_______.png", "Lonar", "+12.3 дБ", SAGE),
    ("S1_VVVH_Wolfe_Creek___________.png", "Wolfe Creek", "+9.3 дБ", GOLD),
]
n9 = 3
gap9 = Inches(0.3)
avail9_in = SLIDE_W.inches - 2 * MARGIN.inches - gap9.inches * (n9 - 1)
col_w9 = Inches(avail9_in / n9)
y9 = Inches(1.85)
img_h9 = Inches(3.15)
for i, (path, name, dbtxt, color) in enumerate(s1_imgs):
    x = MARGIN + i * (col_w9 + gap9)
    ox, oy, w, h = add_image_fit(s, path, x, y9, col_w9, img_h9, center_h=True, center_v=False)
    add_text(s, x, y9 + img_h9 + Inches(0.12), col_w9, Inches(0.4),
              name, size=15, color=CREAM, bold=True, align=PP_ALIGN.CENTER, font=HEAD_FONT)
    add_text(s, x, y9 + img_h9 + Inches(0.55), col_w9, Inches(0.7),
              dbtxt, size=30, color=color, bold=True, align=PP_ALIGN.CENTER, font=HEAD_FONT)
    add_text(s, x, y9 + img_h9 + Inches(1.22), col_w9, Inches(0.3),
              "вал ярче дна (VV)", size=10.5, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
add_text(s, MARGIN, y9 + img_h9 + Inches(1.62), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "Различия обратного рассеяния интерпретируются с учётом шероховатости, влажности, геометрии "
          "наблюдения и морфологии рельефа — это эмпирическое наблюдение по трём объектам, а не общий закон.",
          size=10.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.2)
footer_page(s, next(page_num))

# ---- Слайд 9b: радар vs оптика в литературе -----------------------------------------------
s = add_slide(prs)
add_title(s, "Радар vs оптика для кратеров: что уже сделано",
           "Наш простой VV/VH — базовый вариант давно известного принципа")

col_w9b = Inches(5.75)
gap9b = Inches(0.35)
x1b, x2b = MARGIN, MARGIN + col_w9b + gap9b
y9b = Inches(1.75)
col_h9b = Inches(4.15)

add_card(s, x1b, y9b, col_w9b, col_h9b, color=CARD)
add_text(s, x1b + Inches(0.35), y9b + Inches(0.25), col_w9b - Inches(0.7), Inches(0.4),
          "SIR-C/X-SAR, 1994", size=17, color=TERRACOTTA, bold=True, font=HEAD_FONT)
add_text(s, x1b + Inches(0.35), y9b + Inches(0.68), col_w9b - Inches(0.7), Inches(0.35),
          "McHone et al., 2002, Meteoritics & Planetary Science", size=11, color=MUTED, italic=True)
add_rich(s, x1b + Inches(0.35), y9b + Inches(1.2), col_w9b - Inches(0.7), Inches(2.8),
          [[("10 кратеров радаром, в т.ч. наш Wolfe Creek", CREAM, True, False, 13.5)],
           [("методы: L/C/X-диапазоны, поляризация", MUTED, False, False, 12.5)],
           [("", CREAM, False, False, 6)],
           [("Aorounga (Чад): на оптике — кольцо ~11 км;", CREAM, False, False, 13.5)],
           [("радар вскрыл погребённое под песком второе", CREAM, False, False, 13.5)],
           [("кольцо диаметром >17 км, невидимое в оптике", CREAM, False, False, 13.5)]],
          line_spacing=1.35)

add_card(s, x2b, y9b, col_w9b, col_h9b, color=CARD)
add_text(s, x2b + Inches(0.35), y9b + Inches(0.25), col_w9b - Inches(0.7), Inches(0.4),
          "ALOS PALSAR-2, 2017", size=17, color=SAGE, bold=True, font=HEAD_FONT)
add_text(s, x2b + Inches(0.35), y9b + Inches(0.68), col_w9b - Inches(0.7), Inches(0.35),
          "van Gasselt et al., Earth, Planets and Space", size=11, color=MUTED, italic=True)
add_rich(s, x2b + Inches(0.35), y9b + Inches(1.2), col_w9b - Inches(0.7), Inches(2.8),
          [[("Структура Oasis, Ливия (сильно эродирована)", CREAM, True, False, 13.5)],
           [("методы: L-диапазон, амплитуда, фазовая", MUTED, False, False, 12.5)],
           [("когерентность, поляриметрическая декомпозиция", MUTED, False, False, 12.5)],
           [("", CREAM, False, False, 6)],
           [("Раскрыли подповерхностное напластование и", CREAM, False, False, 13.5)],
           [("палеорусло реки — размытые/невидимые на Landsat", CREAM, False, False, 13.5)]],
          line_spacing=1.35)

y9b_bottom = y9b + col_h9b + Inches(0.2)
add_card(s, MARGIN, y9b_bottom, SLIDE_W - 2 * MARGIN, Inches(1.05), color=CARD)
add_text(s, MARGIN + Inches(0.35), y9b_bottom + Inches(0.14), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(0.8),
          "Общий вывод этой литературы: длинноволновый (L-диапазон) радар проникает сквозь сухой песок "
          "и вскрывает погребённые структуры, невидимые в оптике. Наш анализ (C-диапазон Sentinel-1, "
          "медиана VV/VH) проще и не видит сквозь грунт — это базовый поверхностный сигнал шероховатости "
          "вала; следующий шаг — поляриметрия и L-диапазон (ALOS-2, NISAR).",
          size=11.5, color=CREAM, italic=True, line_spacing=1.22)
footer_page(s, next(page_num))

# ---- Слайд 10: геология (главный результат) -----------------------------------------------
s = add_slide(prs)
add_title(s, "Наши измерения vs публикации vs эталон свежих кратеров")
add_image_fit(s, "crater_geology_comparison.png", MARGIN, Inches(1.55),
               SLIDE_W - 2 * MARGIN, Inches(5.5), center_h=True, center_v=False)
footer_page(s, next(page_num))

# ---- Слайд 11: ограничения -----------------------------------------------
s = add_slide(prs)
add_title(s, "Что метод пока не может")
lim_items = [
    ("Разрешение DEM", "30 м/пиксель не резолвит кратеры мельче ~500 м "
     "(Kaali, 110 м — использованы только опубликованные данные)", TERRACOTTA),
    ("Не заменяет полевую геологию", "определяет морфологию и радарный/рельефный сигнал, "
     "но не шоковый метаморфизм, брекчии, состав пород", GOLD),
    ("Точность справочника координат", "каталожные координаты ориентировочные — "
     "нужна ручная проверка (реализовано: клик по центру на снимке)", SAGE),
    ("Озёра и осадочные заполнения", "искажают измеренную глубину относительно исходной формы "
     "кратера — нужна геологическая интерпретация отклонений, а не только цифра", SALMON),
]
y11 = Inches(1.75)
row_h11 = Inches(1.12)
for i, (head, desc, color) in enumerate(lim_items):
    y = y11 + i * row_h11
    add_circle_num(s, MARGIN, y, Inches(0.55), i + 1, color)
    add_rich(s, MARGIN + Inches(0.85), y - Inches(0.02), Inches(10.9), Inches(0.9),
              [[(head + " — ", color, True, False, 15.5), (desc, CREAM, False, False, 14)]],
              line_spacing=1.2)
footer_page(s, next(page_num))

# ---- Слайд 12: выводы -----------------------------------------------
s = add_slide(prs)
add_title(s, "Выводы и практическая ценность")
concl_items = [
    "Единый воспроизводимый скрипт применён без изменений к трём кратерам разного возраста "
    "и типа мишени (осадочные породы, базальт, песчаник)",
    "Метод валидирован на Барринджере (расхождение с публикацией < 5%); интерпретируемые "
    "расхождения на Lonar и Wolfe Creek объясняются геологией (озеро, занесение осадком)",
    "Радар (Sentinel-1) количественно подтверждает морфологический сигнал вала независимо "
    "от освещения, облачности и типа породы",
    "ДЗЗ — инструмент обнаружения, картографирования и дистанционной проверки согласованности "
    "признаков, но не замена полевой и лабораторной верификации",
]
y12 = Inches(1.75)
for i, txt in enumerate(concl_items):
    y = y12 + i * Inches(1.02)
    add_circle_num(s, MARGIN, y, Inches(0.5), i + 1, TERRACOTTA if i % 2 == 0 else GOLD)
    add_text(s, MARGIN + Inches(0.8), y - Inches(0.03), Inches(11.0), Inches(0.85),
              txt, size=14, color=CREAM, line_spacing=1.2)

add_card(s, MARGIN, Inches(5.95), SLIDE_W - 2 * MARGIN, Inches(1.1), color=CARD)
add_text(s, MARGIN + Inches(0.35), Inches(6.1), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(0.85),
          "Дальше: добавить кратеры с неопределённым происхождением, расширить географию, "
          "автоматизировать первичное детектирование кандидатов.",
          size=12.5, color=GOLD, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
footer_page(s, next(page_num))

# ---- Слайд 13: литература -----------------------------------------------
s = add_slide(prs)
add_title(s, "Литература")
refs_left = [
    "1. Koeberl C. Remote sensing studies of impact craters: how to be sure? // "
    "Comptes Rendus Geoscience. 2004.",
    "2. Gottwald M., et al. Remote sensing of terrestrial impact craters: The TanDEM-X "
    "digital elevation model // Meteoritics & Planetary Science. 2017.",
    "3. Kring D. A. Guidebook to the Geology of Barringer Meteorite Crater, Arizona. "
    "2nd ed. Houston: Lunar and Planetary Institute, 2017.",
]
refs_right = [
    "4. Creation of High Resolution Terrain Models of Barringer Meteorite Crater "
    "(Meteor Crater) Using Photogrammetry and Terrestrial Laser Scanning Methods. "
    "NASA NTRS, 2009.",
    "5. McHone J. F., et al. Space shuttle observations of terrestrial impact structures "
    "using SIR-C and X-SAR radars // Meteoritics & Planetary Science. 2002.",
    "6. van Gasselt S., et al. The Oasis impact structure, Libya: geological "
    "characteristics from ALOS PALSAR-2 data interpretation // "
    "Earth, Planets and Space. 2017.",
]
ref_col_w = Inches((SLIDE_W.inches - 2 * MARGIN.inches - 0.4) / 2)
add_text(s, MARGIN, Inches(1.75), ref_col_w, Inches(3.0),
          "\n\n".join(refs_left), size=12, color=CREAM, line_spacing=1.28)
add_text(s, MARGIN + ref_col_w + Inches(0.4), Inches(1.75), ref_col_w, Inches(3.0),
          "\n\n".join(refs_right), size=12, color=CREAM, line_spacing=1.28)

add_text(s, MARGIN, Inches(5.5), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "Ключевые слова: дистанционное зондирование Земли, импактные структуры, ударные кратеры, "
          "кратер Бэрринджера, Sentinel-1, Sentinel-2, цифровая модель рельефа, радиолокационные данные.",
          size=10.5, color=MUTED, italic=True, line_spacing=1.25)

add_text(s, MARGIN, Inches(6.5), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Спасибо за внимание", size=22, color=TERRACOTTA, bold=True, font=HEAD_FONT,
          align=PP_ALIGN.CENTER)
footer_page(s, next(page_num))

prs.save("crater_presentation.pptx")
print("Сохранено: crater_presentation.pptx")
print(f"Слайдов: {len(prs.slides._sldIdLst)}")
