# -*- coding: utf-8 -*-
"""
Без LibreOffice/PowerPoint под рукой -- QA "на глаз" ненадёжен, поэтому
меряем текст по-настоящему: берём те же самые шрифты (Calibri/Cambria из
C:\\Windows\\Fonts), эмулируем перенос строк по ширине текстового блока и
сравниваем итоговую высоту с высотой блока. Печатает только реальные
переполнения (никаких предупреждений "на всякий случай").
"""

from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

FONT_FILES = {
    ("Calibri", False): "C:/Windows/Fonts/calibri.ttf",
    ("Calibri", True): "C:/Windows/Fonts/calibrib.ttf",
    ("Cambria", False): "C:/Windows/Fonts/cambria.ttc",
    ("Cambria", True): "C:/Windows/Fonts/cambriab.ttf",
}
_font_cache = {}


def get_font(name, bold, size_pt):
    key = (name, bold, round(size_pt * 4))
    if key not in _font_cache:
        path = FONT_FILES.get((name, bold), FONT_FILES[("Calibri", bold)])
        _font_cache[key] = ImageFont.truetype(path, int(size_pt * 4))  # 4x супersampling
    return _font_cache[key]


def text_width_in(text, font, size_pt):
    if not text:
        return 0.0
    px = font.getlength(text)
    # font rendered at size_pt*4 "pixels" (see get_font) -> pt = px/4, in = pt/72
    return (px / 4) / 72.0


def wrap_lines(text, font, size_pt, box_w_in):
    words = text.split(" ")
    lines = []
    cur = ""
    overflow_word = None
    for w in words:
        trial = (cur + " " + w).strip()
        if text_width_in(trial, font, size_pt) <= box_w_in or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
        if text_width_in(w, font, size_pt) > box_w_in:
            overflow_word = w
    if cur:
        lines.append(cur)
    return lines, overflow_word


def check_shape(slide_idx, shape):
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    box_w_in = Emu(shape.width).inches
    box_h_in = Emu(shape.height).inches
    issues = []
    total_h = 0.0
    for p in tf.paragraphs:
        runs = p.runs
        if not runs:
            continue
        full_text = "".join(r.text for r in runs)
        if not full_text.strip():
            total_h += 0.1
            continue
        r0 = runs[0]
        size_pt = r0.font.size.pt if r0.font.size else 14
        bold = bool(r0.font.bold)
        fname = r0.font.name or "Calibri"
        font = get_font(fname, bold, size_pt)
        line_h_in = (size_pt * 1.22) / 72.0
        spacing_mult = p.line_spacing if isinstance(p.line_spacing, (int, float)) else 1.0
        lines, overflow_word = wrap_lines(full_text, font, size_pt, box_w_in)
        n_lines = max(1, len(lines))
        total_h += n_lines * line_h_in * spacing_mult
        if overflow_word:
            issues.append(f"    слово длиннее блока целиком: '{overflow_word[:40]}'")
    if total_h > box_h_in + 0.03:  # небольшой допуск на округления
        issues.append(f"    ПЕРЕПОЛНЕНИЕ: текст ~{total_h:.2f}in > блок {box_h_in:.2f}in "
                       f"(box {box_w_in:.2f}x{box_h_in:.2f}in @ left={Emu(shape.left).inches:.2f} "
                       f"top={Emu(shape.top).inches:.2f})")
    return issues


def check_bounds(slide_idx, shape, slide_w_in, slide_h_in):
    left = Emu(shape.left).inches if shape.left is not None else 0
    top = Emu(shape.top).inches if shape.top is not None else 0
    w = Emu(shape.width).inches if shape.width is not None else 0
    h = Emu(shape.height).inches if shape.height is not None else 0
    issues = []
    if left < -0.01 or top < -0.01 or left + w > slide_w_in + 0.01 or top + h > slide_h_in + 0.01:
        issues.append(f"    ЗА ГРАНИЦЕЙ СЛАЙДА: left={left:.2f} top={top:.2f} "
                       f"right={left+w:.2f} bottom={top+h:.2f} (слайд {slide_w_in:.2f}x{slide_h_in:.2f})")
    return issues


prs = Presentation("crater_presentation.pptx")
slide_w_in = Emu(prs.slide_width).inches
slide_h_in = Emu(prs.slide_height).inches

total_issues = 0
for i, slide in enumerate(prs.slides, 1):
    slide_issues = []
    for shape in slide.shapes:
        text = ""
        if shape.has_text_frame:
            text = " / ".join(p.text for p in shape.text_frame.paragraphs if p.text)[:50]
        issues = check_shape(i, shape) + check_bounds(i, shape, slide_w_in, slide_h_in)
        if issues:
            slide_issues.append((shape.shape_type, text, issues))
    if slide_issues:
        print(f"--- Слайд {i} ---")
        for shape_type, text, issues in slide_issues:
            print(f"  [{shape_type}] \"{text}\"")
            for msg in issues:
                print(msg)
                total_issues += 1

print(f"\nВсего проблем: {total_issues}")
